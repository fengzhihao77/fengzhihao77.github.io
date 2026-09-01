#!/usr/bin/env python3
"""Regenerate images/papers/paper_N.webp from the publication-thumbnail deck.

The deck (images/papers/papers.pptx) is the editable source behind every
publication thumbnail on the site.  It is gitignored, so it only exists on
Zhihao's machine — a fresh clone or CI run will not have it, and this script
exits cleanly in that case rather than failing a build.

The deck's slide canvas is exactly 900:462, the thumbnail aspect ratio, so
"what fills the slide" IS the thumbnail.  There is no framing decision left at
render time: a slide is rendered whole and scaled to 900x462.

Each thumbnail slide carries its target filename in the slide notes (and
nothing else) — that label, not slide position, decides where a slide is
written.  Slides with no label are working material and are skipped.

Render path: LibreOffice --headless --convert-to pdf, then pdftoppm -png at a
DPI chosen so the widest slide lands comfortably above 900px without pointless
upscaling, then cwebp.

Usage:
    python3 scripts/sync_paper_thumbs.py --dry-run   # render, diff, write nothing
    python3 scripts/sync_paper_thumbs.py             # overwrite the thumbnails

Known limitation: slide "paper_1.webp" contains a live PowerPoint chart whose
purple regression line, dashed 95% CI bands and red markers LibreOffice does
not draw.  Regenerating that one slide through this pipeline loses them; export
it from PowerPoint instead.  --dry-run reports the divergence.
"""
from __future__ import annotations

import argparse
import os
import re
import shutil
import subprocess
import sys
import tempfile

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = os.path.join(ROOT, "images", "papers", "papers.pptx")
OUTDIR = os.path.join(ROOT, "images", "papers")

TARGET_W, TARGET_H = 900, 462
WEBP_QUALITY = 72          # reproduces the 23-59 KB range of the committed files
SUPERSAMPLE = 2            # render at 2x, then Lanczos down - kills PDF aliasing
MIN_DPI, MAX_DPI = 96, 600

SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
LABEL_RE = re.compile(r"^paper_\d+\.webp$")

# Expected labels, in index.html top-to-bottom order.
EXPECTED = ["paper_7.webp", "paper_6.webp", "paper_5.webp", "paper_4.webp",
            "paper_3.webp", "paper_2.webp", "paper_1.webp", "paper_0.webp"]


def die(msg: str) -> "None":
    print(f"FAIL  {msg}", file=sys.stderr)
    raise SystemExit(1)


def need(tool: str, hint: str) -> str:
    path = shutil.which(tool)
    if not path:
        die(f"{tool} not found on PATH ({hint})")
    return path


def _pictures(shapes, scale, slide_w):
    """Yield (native_px_wide_used, fraction_of_slide_width) for every picture,
    descending into groups and carrying the group's own scaling."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                inner = (sh.width / sh.element.chExt.cx) if sh.element.chExt.cx else 1.0
            except Exception:
                inner = 1.0
            yield from _pictures(sh.shapes, scale * inner, slide_w)
            continue
        if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            native_w = sh.image.size[0]
            kept = max(1e-6, 1 - (sh.crop_left or 0) - (sh.crop_right or 0))
            yield native_w * kept, (sh.width * scale) / slide_w
        except Exception:
            continue


def read_deck(deck: str) -> "tuple[list[str | None], list[float | None]]":
    """Per slide: the notes label (None = unlabelled) and the worst picture
    resolution ratio (native px available / px the picture occupies at 900 wide).
    A ratio below 1.0 means that picture is upscaled in the final thumbnail.
    None means the slide has no raster content to judge (vector/text/chart)."""
    try:
        from pptx import Presentation
    except ImportError:
        die("python-pptx not installed (pip install python-pptx)")
    prs = Presentation(deck)
    labels, ratios = [], []
    for slide in prs.slides:
        text = ""
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
        labels.append(text if LABEL_RE.match(text) else None)
        worst = None
        for native_w, frac in _pictures(slide.shapes, 1.0, prs.slide_width):
            need_px = max(1.0, frac * TARGET_W)
            r = native_w / need_px
            worst = r if worst is None else min(worst, r)
        ratios.append(worst)
    return labels, ratios


def resolve_mapping(labels: "list[str | None]") -> "list[tuple[int, str]]":
    """Map slide index -> filename, failing loudly on anything ambiguous."""
    labelled = [(i, n) for i, n in enumerate(labels) if n]
    if not labelled:
        print("warn  no slide carries a notes label; falling back to slide order")
        if len(labels) < len(EXPECTED):
            die(f"deck has {len(labels)} slides, need at least {len(EXPECTED)} "
                f"to fall back to slide order")
        return list(zip(range(len(EXPECTED)), EXPECTED))

    names = [n for _, n in labelled]
    dupes = sorted({n for n in names if names.count(n) > 1})
    if dupes:
        die("more than one slide claims the same thumbnail: " + ", ".join(dupes))
    missing = [n for n in EXPECTED if n not in names]
    if missing:
        die("no slide is labelled for: " + ", ".join(missing) +
            "\n      label the slide in its notes with exactly that filename")
    extra = [n for n in names if n not in EXPECTED]
    if extra:
        die("slide labelled for a thumbnail index.html does not use: " + ", ".join(extra))

    order = [n for _, n in labelled]
    if order != EXPECTED:
        print(f"warn  slide order {order} does not match index.html order {EXPECTED}; "
              f"writing by label anyway")
    skipped = len(labels) - len(labelled)
    if skipped:
        print(f"note  {skipped} unlabelled slide(s) skipped (working material)")
    return labelled


def render(deck: str, workdir: str) -> "tuple[list[str], int, tuple[float, float]]":
    """Deck -> one PNG per slide.  Returns (paths, dpi, slide size in points)."""
    need("pdftoppm", "brew install poppler")
    need("pdfinfo", "brew install poppler")
    if not os.path.exists(SOFFICE):
        die(f"LibreOffice not found at {SOFFICE}")

    staged = os.path.join(workdir, "deck.pptx")
    shutil.copy(deck, staged)
    subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf",
                    "--outdir", workdir, staged],
                   check=True, capture_output=True)
    pdf = os.path.join(workdir, "deck.pdf")
    if not os.path.exists(pdf):
        die("LibreOffice produced no PDF")

    info = subprocess.run(["pdfinfo", pdf], check=True, capture_output=True, text=True).stdout
    m = re.search(r"Page size:\s+([\d.]+) x ([\d.]+) pts", info)
    if not m:
        die("could not read page size from pdfinfo")
    pw, ph = float(m.group(1)), float(m.group(2))

    # DPI so the render is at least SUPERSAMPLE x TARGET_W wide, rounded up to a
    # multiple of 24.  Supersampling is about the renderer, not the artwork: it
    # removes PDF rasterisation aliasing and costs nothing but time.
    want = SUPERSAMPLE * TARGET_W * 72 / pw
    dpi = max(MIN_DPI, min(MAX_DPI, int(-(-want // 24) * 24)))
    subprocess.run(["pdftoppm", "-r", str(dpi), "-png", pdf,
                    os.path.join(workdir, "slide")], check=True, capture_output=True)
    pages = sorted((f for f in os.listdir(workdir) if f.startswith("slide-")
                    and f.endswith(".png")),
                   key=lambda f: int(re.findall(r"(\d+)\.png$", f)[0]))
    return [os.path.join(workdir, f) for f in pages], dpi, (pw, ph)


def encode(png: str, dest: str) -> None:
    from PIL import Image
    with Image.open(png) as im:
        im = im.convert("RGB").resize((TARGET_W, TARGET_H), Image.LANCZOS)
        flat = png + ".flat.png"
        im.save(flat)
    subprocess.run(["cwebp", "-quiet", "-q", str(WEBP_QUALITY), flat, "-o", dest],
                   check=True, capture_output=True)


def diff_against_committed(dest: str, committed: str) -> "tuple[str, bool]":
    """Compare a freshly rendered thumbnail with the one currently committed.

    Two numbers, because they answer different questions.  Raw MAE is dominated
    by sub-pixel antialiasing on dense line art and stays around 3-8 even when
    the two images are the same picture.  The same comparison after a light blur
    ignores that edge noise, so it is the one that actually says whether the
    content changed.  The verdict follows the blurred number.
    """
    from PIL import Image, ImageFilter
    if not os.path.exists(committed):
        return "new file (nothing committed at that path)", False
    with Image.open(dest) as a, Image.open(committed) as b:
        size_note = "" if a.size == b.size else f" [size {b.size}->{a.size}]"
        aa = a.convert("RGB").resize((TARGET_W, TARGET_H), Image.LANCZOS)
        bb = b.convert("RGB").resize((TARGET_W, TARGET_H), Image.LANCZOS)
        raw = _mae(aa.tobytes(), bb.tobytes())
        blur = _mae(aa.filter(ImageFilter.GaussianBlur(1.2)).tobytes(),
                    bb.filter(ImageFilter.GaussianBlur(1.2)).tobytes())
    verdict = "identical" if blur < 1.0 else "close" if blur < 4.0 else "DIVERGES"
    kb_new = os.path.getsize(dest) / 1024
    kb_old = os.path.getsize(committed) / 1024
    return (f"MAE={raw:5.2f} struct={blur:5.2f} {verdict:<9s}{size_note} "
            f"{kb_old:.0f}KB->{kb_new:.0f}KB"), verdict == "DIVERGES"


def _mae(a: bytes, b: bytes) -> float:
    return sum(abs(x - y) for x, y in zip(a, b)) / len(a)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--dry-run", action="store_true",
                    help="render to a temp dir and report per-slide diffs; write nothing")
    args = ap.parse_args()

    if not os.path.exists(DECK):
        print(f"note  {os.path.relpath(DECK, ROOT)} is not present — nothing to sync.")
        print("      The deck is gitignored (local-only); this is expected in CI "
              "and fresh clones.")
        return 0

    need("cwebp", "brew install webp")
    try:
        from PIL import Image  # noqa: F401
    except ImportError:
        die("Pillow not installed (pip install Pillow)")

    labels, ratios = read_deck(DECK)
    mapping = resolve_mapping(labels)

    with tempfile.TemporaryDirectory() as work:
        pages, dpi, (pw, ph) = render(DECK, work)
        print(f"note  rendered {len(pages)} slide(s) at {dpi} dpi "
              f"(page {pw:.1f}x{ph:.1f} pt), target {TARGET_W}x{TARGET_H}")
        if len(pages) != len(labels):
            die(f"deck has {len(labels)} slides but the render produced {len(pages)} pages")

        outdir = os.path.join(work, "out") if args.dry_run else OUTDIR
        os.makedirs(outdir, exist_ok=True)

        from PIL import Image
        warned = 0
        diverging = []
        for idx, name in mapping:
            png = pages[idx]
            with Image.open(png) as im:
                sw, sh = im.size
            dest = os.path.join(outdir, name)
            # Snapshot what is committed BEFORE writing.  On a real run dest IS
            # the committed path, so diffing after the write would compare the
            # file with itself and report "identical" no matter what changed.
            baseline = os.path.join(work, "baseline-" + name)
            committed = os.path.join(OUTDIR, name)
            if os.path.exists(committed):
                shutil.copyfile(committed, baseline)
            encode(png, dest)

            # Effective source resolution is set by the artwork on the slide, not
            # by the render DPI - a slide can render at any size and still be a
            # stretched low-res bitmap underneath.
            r = ratios[idx]
            if r is None:
                src = f"src vector {sw}x{sh}"
                flag = ""
            else:
                src = f"src {r:4.2f}x @{sw}x{sh}"
                flag = ""
                if r < 0.98:
                    flag = (f"  WARN slide content only reaches {r:.2f}x the pixels it "
                            f"needs at {TARGET_W}x{TARGET_H} — upscaled, detail is not "
                            f"recoverable from this slide")
                    warned += 1
            detail, diverged = diff_against_committed(dest, baseline)
            if diverged:
                diverging.append(name)
            print(f"  slide{idx + 1:>2} -> {name:<14} {src}  {detail}{flag}")

        if args.dry_run:
            print(f"\nnote  --dry-run: wrote {len(mapping)} file(s) to a temp dir, "
                  f"nothing in {os.path.relpath(OUTDIR, ROOT)} was touched")
        else:
            print(f"\nnote  wrote {len(mapping)} thumbnail(s) to "
                  f"{os.path.relpath(OUTDIR, ROOT)}")
        if warned:
            print(f"warn  {warned} slide(s) were upscaled to reach {TARGET_W}x{TARGET_H}")
        if diverging:
            print(f"warn  content changed vs the committed file: {', '.join(diverging)}")
            if "paper_1.webp" in diverging:
                print("      paper_1.webp is expected: its slide holds a live PowerPoint "
                      "chart whose purple\n      regression line, dashed 95% CI bands and "
                      "red markers LibreOffice does not draw.\n      Export that one slide "
                      "from PowerPoint rather than from this script.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
