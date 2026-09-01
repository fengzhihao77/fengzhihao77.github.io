#!/usr/bin/env python3
"""Regenerate images/talks/talk_N.webp from the conference-thumbnail deck.

The deck (images/talks/talks_v3.pptx) is the editable source behind every
conference thumbnail on the site.  It is gitignored, so it only exists on
Zhihao's machine — a fresh clone or CI run will not have it, and this script
exits cleanly in that case rather than failing a build.

The deck's slide canvas is exactly 4:3 (10in x 7.5in), the thumbnail aspect
ratio, so "what fills the slide" IS the thumbnail.  There is no framing
decision left at render time: a slide is rendered whole and scaled to 700x525.

Each slide holds one full-resolution photo placed edge to edge with a
non-destructive PowerPoint crop (a:srcRect).  Re-framing a photo is therefore
just dragging the crop handles in PowerPoint — and dragging them *outward*
recovers content the current thumbnail cuts off, because the original pixels
are still in the file.

Each thumbnail slide carries its target filename in the slide notes (and
nothing else) — that label, not slide position, decides where a slide is
written.  Slides with no label are working material and are skipped.

Render path: LibreOffice --headless --convert-to pdf, then pdftoppm -png at a
DPI chosen so the slide lands comfortably above 700px without pointless
upscaling, then cwebp.  Rendering the slide (rather than re-implementing
PowerPoint's crop geometry) is what keeps this correct after any edit: crop,
move, resize, rotate, swap the photo, add a second element.

Because the crop is non-destructive, cropping in tighter costs resolution.
Every slide reports the source pixels its current framing actually supplies,
and the headroom left before 700x525 becomes an upscale.

Usage:
    python3 scripts/sync_talk_thumbs.py --dry-run   # render, diff, write nothing
    python3 scripts/sync_talk_thumbs.py             # overwrite the thumbnails
"""
from __future__ import annotations

import argparse
import os
import re
import shutil
import subprocess
import sys
import tempfile

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = os.path.join(ROOT, "images", "talks", "talks_v3.pptx")
OUTDIR = os.path.join(ROOT, "images", "talks")

TARGET_W, TARGET_H = 700, 525
WEBP_QUALITY = 78          # matches the 23-75 KB range of the committed files
# 144 dpi renders the 10in slide at 1440px — a clean ~2x supersample of the
# 700px target.  Downscaling from 2x is measurably cleaner than rasterising
# near the target size (MAE against the committed thumbnails drops by ~1 level
# across the deck), and the photos behind every slide carry far more than
# 1440px, so nothing is being invented.
MIN_DPI, MAX_DPI = 144, 600

SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
LABEL_RE = re.compile(r"^talk_\d+\.webp$")

# Expected labels, in index.html top-to-bottom order.
EXPECTED = ["talk_8.webp", "talk_7.webp", "talk_6.webp", "talk_5.webp",
            "talk_4.webp", "talk_3.webp", "talk_2.webp"]

PICTURE = 13               # MSO_SHAPE_TYPE.PICTURE


def die(msg: str) -> "None":
    print(f"FAIL  {msg}", file=sys.stderr)
    raise SystemExit(1)


def need(tool: str, hint: str) -> str:
    path = shutil.which(tool)
    if not path:
        die(f"{tool} not found on PATH ({hint})")
    return path


def read_deck(deck: str) -> "tuple[list[str | None], list[tuple[int, int] | None]]":
    """Per-slide (notes label, effective source pixels), in slide order.

    "Effective source pixels" is how many original photo pixels the slide's
    framing actually spreads across the full slide — the number that decides
    whether 700x525 is a downscale (good) or an upscale (detail invented).
    It accounts for the crop, for a picture scaled to something other than the
    slide, and for 90/270-degree rotation.  None = no picture on the slide.
    """
    try:
        from pptx import Presentation
    except ImportError:
        die("python-pptx not installed (pip install python-pptx)")

    prs = Presentation(deck)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    labels: "list[str | None]" = []
    source: "list[tuple[int, int] | None]" = []

    for slide in prs.slides:
        text = ""
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
        labels.append(text if LABEL_RE.match(text) else None)

        # The picture covering most of the slide governs its resolution.
        best, best_area = None, 0.0
        for shape in slide.shapes:
            if shape.shape_type != PICTURE:
                continue
            try:
                img_w, img_h = shape.image.size
                cl, ct = shape.crop_left, shape.crop_top
                cr, cb = shape.crop_right, shape.crop_bottom
                sw, sh = shape.width, shape.height
            except Exception:
                continue
            if not sw or not sh:
                continue
            vis_w = img_w * max(0.0, 1.0 - cl - cr)
            vis_h = img_h * max(0.0, 1.0 - ct - cb)
            rot = (shape.rotation or 0.0) % 360
            if 45 <= rot < 135 or 225 <= rot < 315:
                vis_w, vis_h = vis_h, vis_w
                sw, sh = sh, sw
            area = float(sw) * float(sh)
            if area > best_area:
                best_area = area
                best = (int(round(vis_w * slide_w / sw)),
                        int(round(vis_h * slide_h / sh)))
        source.append(best)

    return labels, source


def resolve_mapping(labels: "list[str | None]") -> "list[tuple[int, str]]":
    """Map slide index -> filename, failing loudly on anything ambiguous."""
    labelled = [(i, n) for i, n in enumerate(labels) if n]
    if not labelled:
        print("warn  no slide carries a notes label; falling back to slide order")
        if len(labels) < len(EXPECTED):
            die(f"deck has {len(labels)} slides, need at least {len(EXPECTED)} "
                f"to fall back to slide order")
        return list(zip(range(len(EXPECTED)), EXPECTED))

    names = [n for _, n in labelled]
    dupes = sorted({n for n in names if names.count(n) > 1})
    if dupes:
        die("more than one slide claims the same thumbnail: " + ", ".join(dupes))
    missing = [n for n in EXPECTED if n not in names]
    if missing:
        die("no slide is labelled for: " + ", ".join(missing) +
            "\n      label the slide in its notes with exactly that filename")
    extra = [n for n in names if n not in EXPECTED]
    if extra:
        die("slide labelled for a thumbnail index.html does not use: " + ", ".join(extra))

    order = [n for _, n in labelled]
    if order != EXPECTED:
        print(f"warn  slide order {order} does not match index.html order {EXPECTED}; "
              f"writing by label anyway")
    skipped = len(labels) - len(labelled)
    if skipped:
        print(f"note  {skipped} unlabelled slide(s) skipped (working material)")
    return labelled


def render(deck: str, workdir: str) -> "tuple[list[str], int, tuple[float, float]]":
    """Deck -> one PNG per slide.  Returns (paths, dpi, slide size in points)."""
    need("pdftoppm", "brew install poppler")
    need("pdfinfo", "brew install poppler")
    if not os.path.exists(SOFFICE):
        die(f"LibreOffice not found at {SOFFICE}")

    staged = os.path.join(workdir, "deck.pptx")
    shutil.copy(deck, staged)
    subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf",
                    "--outdir", workdir, staged],
                   check=True, capture_output=True)
    pdf = os.path.join(workdir, "deck.pdf")
    if not os.path.exists(pdf):
        die("LibreOffice produced no PDF")

    info = subprocess.run(["pdfinfo", pdf], check=True, capture_output=True, text=True).stdout
    m = re.search(r"Page size:\s+([\d.]+) x ([\d.]+) pts", info)
    if not m:
        die("could not read page size from pdfinfo")
    pw, ph = float(m.group(1)), float(m.group(2))

    # DPI so the render is at least TARGET_W wide, rounded up to a round number.
    dpi = max(MIN_DPI, min(MAX_DPI, int(-(-TARGET_W * 72 / pw // 24) * 24)))
    subprocess.run(["pdftoppm", "-r", str(dpi), "-png", pdf,
                    os.path.join(workdir, "slide")], check=True, capture_output=True)
    pages = sorted((f for f in os.listdir(workdir) if f.startswith("slide-")
                    and f.endswith(".png")),
                   key=lambda f: int(re.findall(r"(\d+)\.png$", f)[0]))
    return [os.path.join(workdir, f) for f in pages], dpi, (pw, ph)


def encode(png: str, dest: str) -> None:
    from PIL import Image
    with Image.open(png) as im:
        im = im.convert("RGB").resize((TARGET_W, TARGET_H), Image.LANCZOS)
        flat = png + ".flat.png"
        im.save(flat)
    subprocess.run(["cwebp", "-quiet", "-q", str(WEBP_QUALITY), flat, "-o", dest],
                   check=True, capture_output=True)


def diff_against_committed(dest: str, committed: str) -> str:
    from PIL import Image
    import statistics
    if not os.path.exists(committed):
        return "new file (nothing committed at that path)"
    with Image.open(dest) as a, Image.open(committed) as b:
        size_note = "" if a.size == b.size else f" [size {b.size}->{a.size}]"
        aa = a.convert("RGB").resize((TARGET_W, TARGET_H), Image.LANCZOS)
        bb = b.convert("RGB").resize((TARGET_W, TARGET_H), Image.LANCZOS)
        px_a, px_b = aa.tobytes(), bb.tobytes()
    mae = statistics.fmean(abs(x - y) for x, y in zip(px_a, px_b))
    verdict = "identical" if mae < 1 else "close" if mae < 6 else "DIVERGES"
    kb_new = os.path.getsize(dest) / 1024
    kb_old = os.path.getsize(committed) / 1024
    return f"MAE={mae:5.2f} {verdict}{size_note} {kb_old:.0f}KB->{kb_new:.0f}KB"


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--dry-run", action="store_true",
                    help="render to a temp dir and report per-slide diffs; write nothing")
    args = ap.parse_args()

    if not os.path.exists(DECK):
        print(f"note  {os.path.relpath(DECK, ROOT)} is not present — nothing to sync.")
        print("      The deck is gitignored (local-only); this is expected in CI "
              "and fresh clones.")
        return 0

    need("cwebp", "brew install webp")
    try:
        from PIL import Image  # noqa: F401
    except ImportError:
        die("Pillow not installed (pip install Pillow)")

    labels, source = read_deck(DECK)
    mapping = resolve_mapping(labels)

    with tempfile.TemporaryDirectory() as work:
        pages, dpi, (pw, ph) = render(DECK, work)
        print(f"note  rendered {len(pages)} slide(s) at {dpi} dpi "
              f"(page {pw:.1f}x{ph:.1f} pt), target {TARGET_W}x{TARGET_H}")
        if len(pages) != len(labels):
            die(f"deck has {len(labels)} slides but the render produced {len(pages)} pages")

        outdir = os.path.join(work, "out") if args.dry_run else OUTDIR
        os.makedirs(outdir, exist_ok=True)

        warned = 0
        for idx, name in mapping:
            png = pages[idx]
            dest = os.path.join(outdir, name)
            # Snapshot what is committed BEFORE writing.  On a real run dest IS
            # the committed path, so diffing after the write would compare the
            # file with itself and report "identical" no matter what changed.
            baseline = os.path.join(work, "baseline-" + name)
            committed = os.path.join(OUTDIR, name)
            if os.path.exists(committed):
                shutil.copyfile(committed, baseline)
            encode(png, dest)

            eff = source[idx]
            flag = ""
            if eff is None:
                budget = "src n/a (no picture)"
            else:
                sw, sh = eff
                head = min(sw / TARGET_W, sh / TARGET_H)
                budget = f"src {sw}x{sh} headroom {head:4.2f}x"
                if sw < TARGET_W or sh < TARGET_H:
                    flag = (f"  WARN source {sw}x{sh} is below {TARGET_W}x{TARGET_H} — "
                            f"upscaled, detail is not recoverable from this slide")
                    warned += 1
            detail = diff_against_committed(dest, baseline)
            print(f"  slide{idx + 1:>2} -> {name:<14} {budget}  {detail}{flag}")

        if args.dry_run:
            print(f"\nnote  --dry-run: wrote {len(mapping)} file(s) to a temp dir, "
                  f"nothing in {os.path.relpath(OUTDIR, ROOT)} was touched")
        else:
            print(f"\nnote  wrote {len(mapping)} thumbnail(s) to "
                  f"{os.path.relpath(OUTDIR, ROOT)}")
        if warned:
            print(f"warn  {warned} slide(s) were upscaled to reach {TARGET_W}x{TARGET_H}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
